"""
Document processors for multi-format support.

Each processor function takes a file_info tuple and returns (doc_chunks, doc_metadata).
The file_info tuple is:
    (file_path, doc_id, max_chunk_length, chunk_overlap, enable_ocr, ocr_min_width, ocr_min_height, ocr_resolution)
"""

import os
import re
import csv

import pdfplumber
import pytesseract
from PIL import Image

from RAG_Framework.core.text_processing import clean_text
from RAG_Framework.components.indexer.standard import Indexer


def _build_metadata(filename, doc_id, **kwargs):
    """Build a consistent metadata dict for any chunk."""
    meta = {
        'filename': filename,
        'doc_id': doc_id,
        'document_name': filename,
        'full_document_id': f"{filename}_{doc_id}",
    }
    meta.update(kwargs)
    return meta


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def process_pdf(file_info):
    """Process a PDF file — mirrors the original Indexer.process_single_pdf logic."""
    file_path, doc_id, max_chunk_length, chunk_overlap, enable_ocr, ocr_min_width, ocr_min_height, ocr_resolution = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_number, page in enumerate(pdf.pages):
                elements = []

                # Extract tables
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    if table and any(cell for row in table for cell in row if cell):
                        table_text = '\n'.join([' | '.join(map(str, row)) for row in table if row])
                        elements.append({
                            'type': 'table',
                            'content': table_text,
                            'meta': _build_metadata(filename, doc_id, page=page_number, table_idx=table_idx)
                        })

                # Extract paragraph text
                text = page.extract_text()
                if text:
                    text = clean_text(text, source_format='pdf')
                    chunks = Indexer.chunk_text_with_overlap(text, max_chunk_length, chunk_overlap)
                    for chunk_idx, chunk in enumerate(chunks):
                        if chunk:
                            elements.append({
                                'type': 'text',
                                'content': chunk,
                                'meta': _build_metadata(filename, doc_id, page=page_number, chunk_idx=chunk_idx)
                            })

                # OCR for images
                if enable_ocr and page.images:
                    for img_idx, img_dict in enumerate(page.images):
                        try:
                            if 'x0' in img_dict and 'y0' in img_dict and 'x1' in img_dict and 'y1' in img_dict:
                                img_x0 = img_dict["x0"]
                                img_top = img_dict["top"]
                                img_x1 = img_dict["x1"]
                                img_bottom = img_dict["bottom"]

                                if img_x1 <= img_x0 or img_bottom <= img_top:
                                    continue

                                img_x0 = max(0, min(img_x0, page.width))
                                img_x1 = max(0, min(img_x1, page.width))
                                img_top = max(0, min(img_top, page.height))
                                img_bottom = max(0, min(img_bottom, page.height))

                                img_width = img_x1 - img_x0
                                img_height = img_bottom - img_top
                                if img_width < ocr_min_width or img_height < ocr_min_height:
                                    continue

                                image = page.crop((img_x0, img_top, img_x1, img_bottom)).to_image(resolution=ocr_resolution)
                                pil_img = image.original
                                ocr_text = pytesseract.image_to_string(pil_img, lang='eng+por').strip()
                                if ocr_text:
                                    elements.append({
                                        'type': 'image_ocr',
                                        'content': ocr_text,
                                        'meta': _build_metadata(filename, doc_id, page=page_number, img_idx=img_idx)
                                    })
                        except Exception:
                            pass

                for elem in elements:
                    doc_chunks.append(elem['content'])
                    meta = elem['meta'].copy()
                    meta['type'] = elem['type']
                    doc_metadata.append(meta)
        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def process_docx(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, enable_ocr, ocr_min_width, ocr_min_height, ocr_resolution = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        from docx import Document
        doc = Document(file_path)

        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(' | '.join(cells))
            table_text = '\n'.join(rows)
            if table_text.strip():
                doc_chunks.append(table_text)
                meta = _build_metadata(filename, doc_id, page=0, table_idx=table_idx, type='table')
                doc_metadata.append(meta)

        # Extract paragraph text
        full_text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        if full_text:
            full_text = clean_text(full_text, source_format='docx')
            chunks = Indexer.chunk_text_with_overlap(full_text, max_chunk_length, chunk_overlap)
            for chunk_idx, chunk in enumerate(chunks):
                if chunk:
                    doc_chunks.append(chunk)
                    meta = _build_metadata(filename, doc_id, page=0, chunk_idx=chunk_idx, type='text')
                    doc_metadata.append(meta)

        # OCR embedded images
        if enable_ocr:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            img_idx = 0
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_blob = rel.target_part.blob
                        import io
                        pil_img = Image.open(io.BytesIO(image_blob))
                        w, h = pil_img.size
                        if w >= ocr_min_width and h >= ocr_min_height:
                            ocr_text = pytesseract.image_to_string(pil_img, lang='eng+por').strip()
                            if ocr_text:
                                doc_chunks.append(ocr_text)
                                meta = _build_metadata(filename, doc_id, page=0, img_idx=img_idx, type='image_ocr')
                                doc_metadata.append(meta)
                                img_idx += 1
                    except Exception:
                        pass

        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# XLSX / XLS
# ---------------------------------------------------------------------------

def process_xlsx(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, enable_ocr, *_ = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)
        for sheet_index, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                rows.append(' | '.join(cells))
            table_text = '\n'.join(rows)
            if table_text.strip():
                chunks = Indexer.chunk_text_with_overlap(table_text, max_chunk_length, chunk_overlap)
                for chunk_idx, chunk in enumerate(chunks):
                    if chunk:
                        doc_chunks.append(chunk)
                        meta = _build_metadata(
                            filename, doc_id,
                            page=0, sheet=sheet_name, sheet_index=sheet_index,
                            chunk_idx=chunk_idx, type='table'
                        )
                        doc_metadata.append(meta)
        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

def process_csv(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, *_ = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            reader = csv.reader(f)
            rows = []
            for row in reader:
                rows.append(' | '.join(row))
        table_text = '\n'.join(rows)
        if table_text.strip():
            chunks = Indexer.chunk_text_with_overlap(table_text, max_chunk_length, chunk_overlap)
            for chunk_idx, chunk in enumerate(chunks):
                if chunk:
                    doc_chunks.append(chunk)
                    meta = _build_metadata(filename, doc_id, page=0, chunk_idx=chunk_idx, type='table')
                    doc_metadata.append(meta)
        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------

def process_markdown(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, *_ = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()

        # Split by headings to preserve section info
        sections = re.split(r'^(#{1,6}\s+.+)$', raw, flags=re.MULTILINE)

        current_section = ''
        for part in sections:
            part = part.strip()
            if not part:
                continue
            heading_match = re.match(r'^#{1,6}\s+(.+)$', part)
            if heading_match:
                current_section = heading_match.group(1).strip()
                continue

            # Strip markdown syntax
            text = re.sub(r'!\[.*?\]\(.*?\)', '', part)  # images
            text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)  # links → text
            text = re.sub(r'[*_]{1,3}(.+?)[*_]{1,3}', r'\1', text)  # bold/italic
            text = re.sub(r'`{1,3}[^`]*`{1,3}', '', text)  # inline code
            text = re.sub(r'^>\s?', '', text, flags=re.MULTILINE)  # blockquotes
            text = re.sub(r'^[-*+]\s', '', text, flags=re.MULTILINE)  # list markers
            text = re.sub(r'^\d+\.\s', '', text, flags=re.MULTILINE)  # numbered lists

            text = clean_text(text, source_format='markdown')
            if not text:
                continue

            chunks = Indexer.chunk_text_with_overlap(text, max_chunk_length, chunk_overlap)
            for chunk_idx, chunk in enumerate(chunks):
                if chunk:
                    doc_chunks.append(chunk)
                    meta = _build_metadata(
                        filename, doc_id,
                        page=0, section=current_section, chunk_idx=chunk_idx, type='text'
                    )
                    doc_metadata.append(meta)
        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------

def process_plaintext(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, *_ = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                text = f.read()

        text = clean_text(text, source_format='text')
        if text:
            chunks = Indexer.chunk_text_with_overlap(text, max_chunk_length, chunk_overlap)
            for chunk_idx, chunk in enumerate(chunks):
                if chunk:
                    doc_chunks.append(chunk)
                    meta = _build_metadata(filename, doc_id, page=0, chunk_idx=chunk_idx, type='text')
                    doc_metadata.append(meta)
        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def process_pptx(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, enable_ocr, ocr_min_width, ocr_min_height, ocr_resolution = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)

        for slide_idx, slide in enumerate(prs.slides):
            elements = []

            for shape in slide.shapes:
                # Tables
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(' | '.join(cells))
                    table_text = '\n'.join(rows)
                    if table_text.strip():
                        table_idx = len([e for e in elements if e['type'] == 'table'])
                        elements.append({
                            'type': 'table',
                            'content': table_text,
                            'meta': _build_metadata(filename, doc_id, page=slide_idx, slide=slide_idx + 1, table_idx=table_idx)
                        })

                # Text frames
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        text = clean_text(text, source_format='pptx')
                        chunks = Indexer.chunk_text_with_overlap(text, max_chunk_length, chunk_overlap)
                        for chunk_idx, chunk in enumerate(chunks):
                            if chunk:
                                elements.append({
                                    'type': 'text',
                                    'content': chunk,
                                    'meta': _build_metadata(filename, doc_id, page=slide_idx, slide=slide_idx + 1, chunk_idx=chunk_idx)
                                })

                # Images (OCR)
                if enable_ocr and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        import io
                        image_blob = shape.image.blob
                        pil_img = Image.open(io.BytesIO(image_blob))
                        w, h = pil_img.size
                        if w >= ocr_min_width and h >= ocr_min_height:
                            ocr_text = pytesseract.image_to_string(pil_img, lang='eng+por').strip()
                            if ocr_text:
                                img_idx = len([e for e in elements if e['type'] == 'image_ocr'])
                                elements.append({
                                    'type': 'image_ocr',
                                    'content': ocr_text,
                                    'meta': _build_metadata(filename, doc_id, page=slide_idx, slide=slide_idx + 1, img_idx=img_idx)
                                })
                    except Exception:
                        pass

            for elem in elements:
                doc_chunks.append(elem['content'])
                meta = elem['meta'].copy()
                meta['type'] = elem['type']
                doc_metadata.append(meta)

        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def process_html(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, *_ = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        from bs4 import BeautifulSoup

        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            html = f.read()

        soup = BeautifulSoup(html, 'html.parser')

        # Extract tables
        for table_idx, table in enumerate(soup.find_all('table')):
            rows = []
            for tr in table.find_all('tr'):
                cells = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
                rows.append(' | '.join(cells))
            table_text = '\n'.join(rows)
            if table_text.strip():
                doc_chunks.append(table_text)
                meta = _build_metadata(filename, doc_id, page=0, table_idx=table_idx, type='table')
                doc_metadata.append(meta)

        # Extract body text (excluding tables already handled)
        for t in soup.find_all('table'):
            t.decompose()
        body_text = soup.get_text(separator='\n')
        body_text = clean_text(body_text, source_format='html')
        if body_text:
            chunks = Indexer.chunk_text_with_overlap(body_text, max_chunk_length, chunk_overlap)
            for chunk_idx, chunk in enumerate(chunks):
                if chunk:
                    doc_chunks.append(chunk)
                    meta = _build_metadata(filename, doc_id, page=0, chunk_idx=chunk_idx, type='text')
                    doc_metadata.append(meta)

        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# Standalone images (OCR)
# ---------------------------------------------------------------------------

def process_image(file_info):
    file_path, doc_id, max_chunk_length, chunk_overlap, *_ = file_info
    doc_chunks = []
    doc_metadata = []
    filename = os.path.basename(file_path)
    try:
        pil_img = Image.open(file_path)
        ocr_text = pytesseract.image_to_string(pil_img, lang='eng+por').strip()
        ocr_text = clean_text(ocr_text, source_format='image_ocr')
        if ocr_text:
            chunks = Indexer.chunk_text_with_overlap(ocr_text, max_chunk_length, chunk_overlap)
            if chunks:
                for chunk_idx, chunk in enumerate(chunks):
                    if chunk:
                        doc_chunks.append(chunk)
                        meta = _build_metadata(filename, doc_id, page=0, img_idx=chunk_idx, type='image_ocr')
                        doc_metadata.append(meta)
            else:
                doc_chunks.append(ocr_text)
                meta = _build_metadata(filename, doc_id, page=0, img_idx=0, type='image_ocr')
                doc_metadata.append(meta)
        print(f"Processed {filename}")
    except Exception as e:
        print(f"Could not read/process {filename}: {e}")
    return doc_chunks, doc_metadata


# ---------------------------------------------------------------------------
# Extension → processor mapping
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    '.pdf': process_pdf,
    '.docx': process_docx,
    '.xlsx': process_xlsx,
    '.xls': process_xlsx,
    '.csv': process_csv,
    '.md': process_markdown,
    '.txt': process_plaintext,
    '.pptx': process_pptx,
    '.html': process_html,
    '.htm': process_html,
    '.png': process_image,
    '.jpg': process_image,
    '.jpeg': process_image,
    '.tiff': process_image,
    '.bmp': process_image,
}


def process_single_document(file_info):
    """
    Top-level dispatch function (picklable for multiprocessing).
    Looks up the correct processor based on file extension.
    """
    file_path = file_info[0]
    ext = os.path.splitext(file_path)[1].lower()
    processor = SUPPORTED_EXTENSIONS.get(ext)
    if processor is None:
        filename = os.path.basename(file_path)
        print(f"Unsupported file format: {filename} ({ext})")
        return [], []
    return processor(file_info)
